import os
import json
import traceback
import requests
from datetime import datetime
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# === Google Drive Setup ===
drive_service = None
try:
    creds_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if creds_json:
        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/drive"]
        )
        drive_service = build("drive", "v3", credentials=creds)
except Exception as e:
    print(f"❌ Drive setup failed: {e}")
    traceback.print_exc()

def upload_to_drive(file_path, session_id):
    try:
        query = f"name='{session_id}' and mimeType='application/vnd.google-apps.folder'"
        results = drive_service.files().list(q=query, fields="files(id)").execute()
        folders = results.get("files", [])
        if folders:
            folder_id = folders[0]["id"]
        else:
            folder = drive_service.files().create(body={
                "name": session_id,
                "mimeType": "application/vnd.google-apps.folder"
            }, fields="id").execute()
            folder_id = folder["id"]

        file_meta = {"name": os.path.basename(file_path), "parents": [folder_id]}
        media = MediaFileUpload(file_path, resumable=True)
        uploaded = drive_service.files().create(body=file_meta, media_body=media, fields="id").execute()
        return f"https://drive.google.com/file/d/{uploaded['id']}/view"
    except Exception as e:
        print(f"❌ Upload failed: {e}")
        traceback.print_exc()
        return None

def download_files(files, folder_path):
    downloaded = []
    for f in files:
        if not f.get("file_url"):
            print(f"⚠️ Skipping {f.get('file_name')} – no file_url")
            continue
        try:
            path = os.path.join(folder_path, f["file_name"])
            r = requests.get(f["file_url"], timeout=20)
            with open(path, "wb") as out:
                out.write(r.content)
            f["local_path"] = path
            downloaded.append(f)
        except Exception as e:
            print(f"❌ Download failed for {f['file_name']}: {e}")
    return downloaded

def generate_financial_excel(session_id, folder_path):
    path = os.path.join(folder_path, f"IT_Financial_Analysis_{session_id}.xlsx")
    wb = Workbook()

    # === Sheet 1: Executive Summary ===
    ws = wb.active
    ws.title = "Executive Summary"
    ws.append(["Metric", "Value"])
    ws.append(["Total Project Cost", "$3,500,000"])
    ws.append(["Estimated Annual Savings", "$950,000"])
    ws.append(["Payback Period (Months)", "22"])
    ws.append(["ROI (%)", "68.2"])
    ws.append(["NPV @7%", "$1,200,000"])
    ws.append(["NPV @10%", "$950,000"])
    ws.append(["Breakeven Date", "Month 23"])

    # === Additional Sheets Placeholder ===
    for title in [
        "Current vs Target Cost",
        "Phased Investment Plan",
        "Savings Forecast (5Y)",
        "Risk Exposure Cost",
        "TCO Breakdown",
        "Workforce Impact",
        "ROI/NPV Scenarios",
        "Benchmarks",
        "Assumptions & Notes"
    ]:
        ws = wb.create_sheet(title)
        ws.append([f"{title} (Data Placeholder)"])

    wb.save(path)
    return path

def generate_financial_pptx(session_id, folder_path):
    path = os.path.join(folder_path, f"IT_Financial_Summary_Deck_{session_id}.pptx")
    ppt = Presentation()
    
    def add_slide(title, bullets):
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for bullet in bullets:
            tf.add_paragraph().text = bullet

    # === Slides ===
    add_slide("Executive Summary", [
        "Total Investment: $3.5M",
        "ROI: 68.2%",
        "Breakeven: Month 23"
    ])
    add_slide("TCO Breakdown", [
        "Hardware: 35%",
        "Cloud: 25%",
        "Services: 20%",
        "Support: 10%",
        "Other: 10%"
    ])
    add_slide("Savings Forecast", [
        "Year 1: $200,000",
        "Year 2: $400,000",
        "Year 3: $650,000",
        "Year 4: $850,000",
        "Year 5: $950,000"
    ])
    add_slide("Risk Exposure", [
        "Non-compliance Risk: $150K",
        "Downtime Exposure: $120K",
        "Cost Overrun Buffer: $200K"
    ])
    add_slide("Phase-wise Investment", [
        "Phase 1: $1.2M",
        "Phase 2: $1.0M",
        "Phase 3: $900K",
        "Phase 4: $400K"
    ])

    ppt.save(path)
    return path

def generate_financial_docx(session_id, folder_path):
    path = os.path.join(folder_path, f"IT_Investment_Brief_{session_id}.docx")
    doc = Document()
    doc.add_heading("IT Investment Brief", 0)
    doc.add_paragraph(f"Session: {session_id}\n")

    doc.add_heading("1. Financial Justification", level=1)
    doc.add_paragraph("The total projected investment of $3.5M will yield an estimated annual saving of $950K and a payback period of less than 24 months.")

    doc.add_heading("2. Strategic Alignment", level=1)
    doc.add_paragraph("This financial plan aligns with our modernization roadmap and compliance posture, optimizing cost while enabling agility.")

    doc.add_heading("3. Funding Recommendation", level=1)
    doc.add_paragraph("We recommend phased funding approval with milestone-based release to mitigate risk and ensure alignment.")

    doc.save(path)
    return path

def process_financials(session_id, email, files, folder_path):
    try:
        os.makedirs(folder_path, exist_ok=True)
        downloaded = download_files(files, folder_path)

        xlsx_path = generate_financial_excel(session_id, folder_path)
        pptx_path = generate_financial_pptx(session_id, folder_path)
        docx_path = generate_financial_docx(session_id, folder_path)

        xlsx_url = upload_to_drive(xlsx_path, session_id)
        pptx_url = upload_to_drive(pptx_path, session_id)
        docx_url = upload_to_drive(docx_path, session_id)

        for f in downloaded:
            f["file_url"] = upload_to_drive(f["local_path"], session_id)

        downloaded.extend([
            {
                "file_name": os.path.basename(xlsx_path),
                "file_url": xlsx_url,
                "file_type": "xlsx_financial"
            },
            {
                "file_name": os.path.basename(pptx_path),
                "file_url": pptx_url,
                "file_type": "pptx_financial"
            },
            {
                "file_name": os.path.basename(docx_path),
                "file_url": docx_url,
                "file_type": "docx_financial_brief"
            }
        ])

        next_gpt_url = "https://it-summarizer-api.onrender.com/start_summarizer"
        payload = {
            "session_id": session_id,
            "email": email,
            "gpt_module": "it_financials",
            "files": downloaded,
            "status": "complete"
        }

        requests.post(next_gpt_url, json=payload)

    except Exception as e:
        print(f"🔥 Financial processing failed: {e}")
        traceback.print_exc()
